import sys
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

# 한글 폰트 설정 (Windows 기본 맑은 고딕)
try:
    pdfmetrics.registerFont(TTFont('Malgun', 'malgun.ttf'))
    FONT_NAME = 'Malgun'
except Exception:
    FONT_NAME = 'Helvetica'  # Fallback (한글 깨짐 발생 가능)

def format_price(amount):
    """숫자를 'O억 O천만원' 형식으로 가독성 있게 포맷팅"""
    try:
        val = int(float(amount))
        if val == 0: return "0원"
        uk = val // 100000000
        man = (val % 100000000) // 10000
        res = ""
        if uk > 0: res += f"{uk}억 "
        if man > 0: res += f"{man}만 "
        res += "원"
        return res.strip()
    except Exception:
        return str(amount)

def generate_pptx_bytes(rows: list[dict]) -> bytes:
    """영업 및 설득용 PPTX 바이트 스트림 생성"""
    prs = Presentation()
    for r in rows:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = f"경매취하 솔루션 제안 [{r.get('사건번호', '사건미상')}]"
        
        # 우측 상단 회사 로고 자리 (Placeholder)
        shape = slide.shapes.add_shape(1, Inches(7.5), Inches(0.2), Inches(2.0), Inches(0.5))
        shape.text = "회사 로고 입력란"
        
        body = slide.shapes.placeholders[1].text_frame
        body.text = f"목적물: {r.get('주소', '')} {r.get('아파트명', '')}\n"
        
        p = body.add_paragraph()
        p.text = "■ 자금 및 가치 현황"
        p = body.add_paragraph()
        p.level = 1
        p.text = f"감정가: {format_price(r.get('감정가'))} / 예상가: {format_price(r.get('낙찰예상가'))}"
        p = body.add_paragraph()
        p.level = 1
        p.text = f"부채총액(청구액 등): {format_price(r.get('부채총액'))}"
        
        p = body.add_paragraph()
        p.text = "■ 권리 분석 서머리"
        p = body.add_paragraph()
        p.level = 1
        p.text = str(r.get("권리요약", "특이사항 없음"))
        
        p = body.add_paragraph()
        p.text = "■ 종합 제안 (소유주/채권자 설득 전략)"
        
        memo_lines = str(r.get("담당자메모", "상담 진행 권장")).split('\n')
        for line in memo_lines:
            if line.strip():
                p = body.add_paragraph()
                p.level = 1
                p.text = line.strip()

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()

def generate_pdf_bytes(rows: list[dict]) -> bytes:
    """보고용 PDF 바이트 스트림 생성"""
    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=A4)
    width, height = A4
    for r in rows:
        y = height - 60
        c.setFont(FONT_NAME, 22)
        c.drawString(50, y, "경매취하 및 자금조달 검토 보고서")
        
        y -= 40
        c.setFont(FONT_NAME, 14)
        c.drawString(50, y, f"사건번호: {r.get('사건번호', '사건미상')}")
        
        # 기본 정보 섹션
        y -= 40
        c.setFont(FONT_NAME, 12)
        c.drawString(50, y, "1. 물건 요약")
        y -= 20
        c.setFont(FONT_NAME, 10)
        c.drawString(60, y, f"소 재 지: {r.get('주소', '')} {r.get('아파트명', '')}")
        
        # 금액 분석 섹션
        y -= 35
        c.setFont(FONT_NAME, 12)
        c.drawString(50, y, "2. 자금 및 권리 현황")
        y -= 20
        c.setFont(FONT_NAME, 10)
        
        fields = [
            ("감 정 가", format_price(r.get('감정가'))),
            ("예 상 가", format_price(r.get('낙찰예상가'))),
            ("부채총액", format_price(r.get('부채총액'))),
            ("분석점수", str(r.get('분석점수', '')) + " 점 (" + str(r.get('분석등급', '')) + "등급)"),
            ("권리현황", str(r.get('권리요약', '이슈 없음'))),
        ]
        
        for k, v in fields:
            c.drawString(60, y, f"- {k} : {v}")
            y -= 18
        
        # 제안 및 멘트 섹션
        y -= 20
        c.setFont(FONT_NAME, 12)
        c.drawString(50, y, "3. 대응 솔루션 및 상담 가이드")
        y -= 20
        c.setFont(FONT_NAME, 10)
        
        memo_lines = str(r.get("담당자메모", "")).split('\n')
        for line in memo_lines:
            if line.strip():
                # 간단한 자동 줄바꿈 대신 그대로 인쇄 (MVP용)
                c.drawString(60, y, line.strip())
                y -= 18
        
        # 하단 푸터 (로고 대용)
        c.setFont(FONT_NAME, 9)
        c.drawString(width - 150, 40, "회사 로고 위치 / 대외비 문구")
        c.showPage()
    
    c.save()
    buf.seek(0)
    return buf.getvalue()
